from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
import PyPDF2
import pdfplumber
from docx import Document
import openpyxl
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_bytes
from PIL import Image
import io
import base64
from langdetect import detect
import tempfile
import os
from typing import Dict, Any, List

app = FastAPI(title="Document Extractor Service")

@app.get("/")
async def health_check():
    return {"status": "healthy", "service": "document-extractor"}

def create_smart_chunks(text: str, chunk_size: int = 4000, overlap: int = 400) -> List[Dict[str, Any]]:
    """Crea chunks inteligentes del texto"""
    if not text or len(text) <= chunk_size:
        return [{
            'content': text,
            'chunk_index': 0,
            'total_chunks': 1,
            'start_position': 0,
            'end_position': len(text) if text else 0
        }]
    
    chunks = []
    
    # Dividir por párrafos primero (más inteligente que por oraciones)
    paragraphs = text.split('\n\n')
    current_chunk = ""
    chunk_index = 0
    start_position = 0
    
    for paragraph in paragraphs:
        # Si el párrafo solo cabe en un nuevo chunk
        if len(current_chunk) + len(paragraph) + 2 > chunk_size:
            if current_chunk:
                chunks.append({
                    'content': current_chunk.strip(),
                    'chunk_index': chunk_index,
                    'start_position': start_position,
                    'end_position': start_position + len(current_chunk)
                })
                chunk_index += 1
                
                # Overlap: tomar las últimas líneas del chunk anterior
                lines = current_chunk.split('\n')
                overlap_lines = lines[-5:] if len(lines) > 5 else lines
                overlap_text = '\n'.join(overlap_lines)
                
                start_position = start_position + len(current_chunk) - len(overlap_text)
                current_chunk = overlap_text + "\n\n" + paragraph
            else:
                # Si un párrafo es muy largo, dividirlo por oraciones
                if len(paragraph) > chunk_size:
                    sentences = paragraph.split('. ')
                    for sentence in sentences:
                        if len(current_chunk) + len(sentence) + 2 > chunk_size:
                            chunks.append({
                                'content': current_chunk.strip(),
                                'chunk_index': chunk_index,
                                'start_position': start_position,
                                'end_position': start_position + len(current_chunk)
                            })
                            chunk_index += 1
                            start_position += len(current_chunk)
                            current_chunk = sentence + ". "
                        else:
                            current_chunk += sentence + ". "
                else:
                    current_chunk = paragraph
        else:
            if current_chunk:
                current_chunk += "\n\n" + paragraph
            else:
                current_chunk = paragraph
    
    # Agregar el último chunk
    if current_chunk:
        chunks.append({
            'content': current_chunk.strip(),
            'chunk_index': chunk_index,
            'start_position': start_position,
            'end_position': len(text)
        })
    
    # Actualizar total_chunks
    for chunk in chunks:
        chunk['total_chunks'] = len(chunks)
    
    return chunks

@app.post("/extract")
async def extract_document(
    file: UploadFile = File(...),
    extract_tables: bool = True,
    detect_language: bool = True,
    ocr_when_needed: bool = True,
    chunk_text: bool = True,
    chunk_size: int = 4000,
    chunk_overlap: int = 400
):
    """
    Extrae texto de varios formatos de documentos con opción de chunking
    """
    try:
        # Leer el archivo
        content = await file.read()
        filename = file.filename.lower()
        
        # Determinar el tipo de archivo
        if filename.endswith('.pdf'):
            result = await extract_pdf(content, extract_tables, ocr_when_needed)
        elif filename.endswith(('.docx', '.doc')):
            result = await extract_docx(content)
        elif filename.endswith(('.xlsx', '.xls')):
            result = await extract_excel(content)
        elif filename.endswith(('.pptx', '.ppt')):
            result = await extract_pptx(content)
        elif filename.endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp')):
            result = await extract_image(content)
        elif filename.endswith(('.txt', '.md')):
            result = await extract_text(content)
        else:
            raise HTTPException(status_code=400, detail=f"Formato no soportado: {filename}")
        
        # Detectar idioma si se solicita
        if detect_language and result.get('text'):
            try:
                result['language'] = detect(result['text'][:500])
            except:
                result['language'] = 'unknown'
        
        # Agregar metadata
        result['filename'] = file.filename
        result['file_size'] = len(content)
        result['content_type'] = file.content_type
        
        # Crear chunks si se solicita
        if chunk_text and result.get('text'):
            result['chunks'] = create_smart_chunks(result['text'], chunk_size, chunk_overlap)
            result['total_chunks'] = len(result['chunks'])
            result['chunking_method'] = 'smart_paragraph_based'
            result['chunk_size'] = chunk_size
            result['chunk_overlap'] = chunk_overlap
        
        return result
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

async def extract_pdf(content: bytes, extract_tables: bool = True, use_ocr: bool = True) -> Dict[str, Any]:
    """Extrae texto de PDF con soporte para OCR y tablas"""
    result = {
        'text': '',
        'tables': [],
        'pages': 0,
        'needs_ocr': False,
        'metadata': {}
    }
    
    # Intentar extracción normal primero
    try:
        pdf_file = io.BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        result['pages'] = len(pdf_reader.pages)
        
        text_parts = []
        for page_num, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(f"--- Página {page_num + 1} ---\n{page_text}")
            else:
                result['needs_ocr'] = True
        
        result['text'] = '\n\n'.join(text_parts)
        
        # Si no hay texto o muy poco, intentar OCR
        if (not result['text'].strip() or len(result['text']) < 100) and use_ocr:
            result['needs_ocr'] = True
            result['text'] = await extract_pdf_with_ocr(content)
        
        # Extraer tablas si se solicita
        if extract_tables:
            result['tables'] = await extract_pdf_tables(content)
            
    except Exception as e:
        # Si falla PyPDF2, intentar con pdfplumber
        try:
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_file:
                tmp_file.write(content)
                tmp_path = tmp_file.name
                
            with pdfplumber.open(tmp_path) as pdf:
                result['pages'] = len(pdf.pages)
                text_parts = []
                
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Página {i + 1} ---\n{page_text}")
                
                result['text'] = '\n\n'.join(text_parts)
                
            os.unlink(tmp_path)
        except:
            if use_ocr:
                result['text'] = await extract_pdf_with_ocr(content)
                result['needs_ocr'] = True
    
    return result

async def extract_pdf_with_ocr(content: bytes) -> str:
    """Extrae texto de PDF usando OCR"""
    try:
        images = convert_from_bytes(content, dpi=300)
        text_parts = []
        
        for i, image in enumerate(images):
            text = pytesseract.image_to_string(image, lang='spa+eng')
            text_parts.append(f"--- Página {i + 1} (OCR) ---\n{text}")
        
        return '\n\n'.join(text_parts)
    except Exception as e:
        return f"Error en OCR: {str(e)}"

async def extract_pdf_tables(content: bytes) -> List[Dict[str, Any]]:
    """Extrae tablas de PDF - temporalmente simplificado"""
    tables_data = []
    
    try:
        # Por ahora usamos pdfplumber para extraer tablas básicas
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name
        
        with pdfplumber.open(tmp_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                for i, table in enumerate(tables):
                    if table:
                        tables_data.append({
                            'table_index': i,
                            'page': page_num + 1,
                            'data': table
                        })
        
        os.unlink(tmp_path)
    except Exception as e:
        # Si falla, simplemente retornar lista vacía
        pass
    
    return tables_data

async def extract_docx(content: bytes) -> Dict[str, Any]:
    """Extrae texto de archivos Word"""
    try:
        doc_stream = io.BytesIO(content)
        doc = Document(doc_stream)
        
        text_parts = []
        tables_data = []
        
        # Extraer párrafos
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extraer tablas
        for i, table in enumerate(doc.tables):
            table_content = []
            for row in table.rows:
                row_content = [cell.text for cell in row.cells]
                table_content.append(row_content)
            
            tables_data.append({
                'table_index': i,
                'data': table_content
            })
        
        return {
            'text': '\n\n'.join(text_parts),
            'tables': tables_data,
            'paragraphs': len(doc.paragraphs)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error procesando DOCX: {str(e)}")

async def extract_excel(content: bytes) -> Dict[str, Any]:
    """Extrae datos de archivos Excel"""
    try:
        excel_stream = io.BytesIO(content)
        workbook = openpyxl.load_workbook(excel_stream, data_only=True)
        
        sheets_data = []
        all_text = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_content = []
            
            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        row_values.append(str(cell.value))
                
                if row_values:
                    sheet_content.append(row_values)
                    all_text.append(' | '.join(row_values))
            
            if sheet_content:
                sheets_data.append({
                    'sheet_name': sheet_name,
                    'data': sheet_content,
                    'rows': sheet.max_row,
                    'columns': sheet.max_column
                })
        
        return {
            'text': '\n'.join(all_text),
            'sheets': sheets_data,
            'total_sheets': len(workbook.sheetnames)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error procesando Excel: {str(e)}")

async def extract_pptx(content: bytes) -> Dict[str, Any]:
    """Extrae texto de presentaciones PowerPoint"""
    try:
        pptx_stream = io.BytesIO(content)
        presentation = Presentation(pptx_stream)
        
        slides_content = []
        all_text = []
        
        for i, slide in enumerate(presentation.slides):
            slide_text = []
            
            # Extraer texto de shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            # Extraer notas del presentador
            if slide.notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    slide_text.append(f"[Notas: {notes_text}]")
            
            slide_content = '\n'.join(slide_text)
            if slide_content:
                slides_content.append({
                    'slide_number': i + 1,
                    'content': slide_content
                })
                all_text.append(f"--- Diapositiva {i + 1} ---\n{slide_content}")
        
        return {
            'text': '\n\n'.join(all_text),
            'slides': slides_content,
            'total_slides': len(presentation.slides)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error procesando PowerPoint: {str(e)}")

async def extract_image(content: bytes) -> Dict[str, Any]:
    """Extrae texto de imágenes usando OCR"""
    try:
        image = Image.open(io.BytesIO(content))
        
        # OCR
        text = pytesseract.image_to_string(image, lang='spa+eng')
        
        # Metadata de la imagen
        metadata = {
            'format': image.format,
            'mode': image.mode,
            'size': image.size,
            'width': image.width,
            'height': image.height
        }
        
        return {
            'text': text.strip(),
            'metadata': metadata,
            'has_text': bool(text.strip())
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error procesando imagen: {str(e)}")

async def extract_text(content: bytes) -> Dict[str, Any]:
    """Extrae texto de archivos de texto plano"""
    try:
        text = content.decode('utf-8')
        return {
            'text': text,
            'lines': len(text.split('\n')),
            'characters': len(text)
        }
    except UnicodeDecodeError:
        # Intentar con otro encoding
        text = content.decode('latin-1')
        return {
            'text': text,
            'lines': len(text.split('\n')),
            'characters': len(text),
            'encoding': 'latin-1'
        }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
